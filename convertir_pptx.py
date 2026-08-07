"""Convierte PPTX a Markdown: texto por diapositiva, tablas e imágenes extraídas.

Uso (ejecutar desde el directorio que contiene los pptx):
    python convertir_pptx.py "carpeta/presentacion.pptx" ["otra.pptx" ...]

Salida: output/<misma estructura>/<presentacion>/<presentacion>.md
        con imágenes en .../assets/.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path.cwd()


def slugify(name: str) -> str:
    return name.replace("(", "-").replace(")", "-").replace(" ", "_")


def extract_text(shape):
    if shape.has_text_frame:
        lines = []
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs)
            if text.strip():
                lines.append(text)
        return lines
    return []


def convert_pptx(pptx_path: Path):
    rel = pptx_path.resolve().relative_to(ROOT)
    doc_name = pptx_path.stem
    out_dir = ROOT / "output" / slugify(str(rel.parent)) / slugify(doc_name)
    assets_dir = out_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    md_lines = [f"# {doc_name}\n"]
    img_count = 0

    for i, slide in enumerate(prs.slides, start=1):
        md_lines.append(f"## Diapositiva {i}\n")

        title = None
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                if title:
                    md_lines.append(f"**{title}**\n")
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_count += 1
                ext = shape.image.ext
                img_name = f"slide{i:03d}-{img_count:03d}.{ext}"
                (assets_dir / img_name).write_bytes(shape.image.blob)
                md_lines.append(f"![](assets/{img_name})\n")
            elif shape.has_table:
                table = shape.table
                rows = [[cell.text for cell in row.cells] for row in table.rows]
                if rows:
                    md_lines.append("| " + " | ".join(rows[0]) + " |")
                    md_lines.append("|" + "---|" * len(rows[0]))
                    for r in rows[1:]:
                        md_lines.append("| " + " | ".join(r) + " |")
                    md_lines.append("")
            else:
                lines = extract_text(shape)
                for line in lines:
                    md_lines.append(f"- {line}")
                if lines:
                    md_lines.append("")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md_lines.append(f"> Notas: {notes}\n")

        md_lines.append("")

    md_text = "\n".join(md_lines)
    md_path = out_dir / f"{doc_name}.md"
    md_path.write_text(md_text, encoding="utf-8")
    return md_path, len(prs.slides), img_count


def main(paths):
    if not paths:
        print("Uso: python convertir_pptx.py archivo1.pptx [archivo2.pptx ...]", file=sys.stderr)
        sys.exit(1)
    for p in paths:
        pptx_path = Path(p)
        print(f"Convirtiendo: {pptx_path.name} ...")
        md_path, n_slides, n_img = convert_pptx(pptx_path)
        print(f"  OK -> {md_path.relative_to(ROOT)} | {n_slides} diapositivas | {n_img} img")


if __name__ == "__main__":
    main(sys.argv[1:])
